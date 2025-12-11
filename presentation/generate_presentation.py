from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {
        "title": "BlazorGameEngine - Architettura Interna",
        "bullets": [
            "Engine 2D per Blazor + HTML5 Canvas",
            "API ispirata a GameMaker, type-safe in C#",
            "Progettato per performance: batching, DI, async"
        ],
        "notes": "Introduzione rapida: obiettivo della presentazione e target audience: sviluppatori C# / Blazor interessati a game engine 2D."
    },
    {
        "title": "Agenda",
        "bullets": [
            "Architettura generale",
            "Game loop e interop JavaScript",
            "Gestione input, rendering e fisica",
            "Dependency Injection per asset e GameObject",
            "Performance, limitazioni e roadmap"
        ],
        "notes": "Dire il percorso della presentazione e quanto tempo per demo/domande."
    },
    {
        "title": "Architettura High-level",
        "bullets": [
            "Blazor Component (GameView) ↔ Game (C#) ↔ Canvas2D (JS)",
            "Game loop gestito da JavaScript con requestAnimationFrame",
            "Batching Canvas2D per ridurre l'overhead di interop"
        ],
        "notes": "Mostrare diapositiva/diagramma a blocchi. Evidenziare responsabilità di ogni layer."
    },
    {
        "title": "Game Loop & JS Interop",
        "bullets": [
            "Loop: JS invoke C# frame handler + requestAnimationFrame",
            "Limitazione FPS con setTimeout per target rate",
            "Sequenza: render -> update -> input state snapshot"
        ],
        "notes": "Spiegare perché il loop è in JS: timing preciso e sincronizzazione col browser."
    },
    {
        "title": "Esempio: ExecuteFrameAsync (sequenza)",
        "bullets": [
            "Aggiorna mouse, snapshot istanze",
            "Fase di rendering: OnDrawAsync + GameObject.OnDrawAsync",
            "Fase di update: OnStepAsync + GameObject.ExecuteStepAsync",
            "Salva stato input e aggiorna controllerHub"
        ],
        "notes": "Mostrare snippet semplificato (solo pseudocodice) e spiegare uso di ToArray() per evitare concorrenza."
    },
    {
        "title": "Gestione Input: Pattern 'Previous State'",
        "bullets": [
            "Doppio buffer: current vs previous per tasti e mouse",
            "Check: hold | CheckPressed: edge detect | CheckReleased",
            "Permette logiche frame-accurate (es. single-tap)"
        ],
        "notes": "Mostrare tabella esemplificativa di frame con space key. Sottolineare semplicità e robustezza."
    },
    {
        "title": "GameObject & Dependency Injection",
        "bullets": [
            "Istanzia con ActivatorUtilities + reflection per init-only",
            "GameObject riceve dipendenze via DI (asset, servizi)",
            "OnCreate callback e proprietà immutabili OriginalX/OriginalY"
        ],
        "notes": "Spiegare il vantaggio: testabilità, iniezione di asset e servizi condivisi."
    },
    {
        "title": "Rendering: Sprite, Transform, Batching",
        "bullets": [
            "Distinzione ImageAsset vs SpriteAsset (sheet frames)",
            "DrawSprite calcola imageX/imageY; supporto scale/rotate",
            "BeginBatch/EndBatch per minimizzare round-trips JS↔C#"
        ],
        "notes": "Mostrare esempio di calcolo frame e ricordare l'ottimizzazione con batching."
    },
    {
        "title": "Fisica e Collisioni",
        "bullets": [
            "AABB per collision detection (IsCollidingWith)",
            "MoveContactSolid: movimento con stop al contatto (pixel-step)",
            "MoveOutsideSolid per espulsione da overlapping"
        ],
        "notes": "Discutere trade-off: semplicità vs precisione (no pixel-perfect). Possibili ottimizzazioni: spatial hashing."
    },
    {
        "title": "Asset Management & Registrazione DI",
        "bullets": [
            "Estensione: AddSprite<T> registra come T, SpriteAsset, IGameAsset",
            "Auto-discovery con reflection sull'assembly",
            "Consente enumerazione e iniezione forte-typed"
        ],
        "notes": "Sottolineare pattern per scalabilità e facilità d'uso per game devs."
    },
    {
        "title": "Tipo Angle e Safety",
        "bullets": [
            "Struct Angle con factory FromDegrees/FromRadians",
            "Metodi Sin/Cos e operatori per composizione",
            "Evita ambiguità gradi vs radianti nel codice" 
        ],
        "notes": "Mostrare breve snippet. Spiegare beneficio in refactoring e leggibilità."
    },
    {
        "title": "Performance & Ottimizzazioni",
        "bullets": [
            "Ridurre overhead interop con batching",
            "Usare snapshot ToArray per evitare lock durante iterazioni",
            "Possibili futuri: WebGL, object pooling, spatial hashing"
        ],
        "notes": "Dare numeri qualitativi: quanti draw calls si possono ridurre con batching."
    },
    {
        "title": "Limitazioni e Roadmap",
        "bullets": [
            "Single-threaded (Blazor WASM)",
            "No ECS, collision AABB solo",
            "Evoluzioni: WebGL, pooling, streaming asset"
        ],
        "notes": "Discutere trade-offs e possibili PR/feature requests."
    },
    {
        "title": "Demo / Come eseguire il gioco",
        "bullets": [
            "Aprire soluzione `ChristmasJumpGame.slnx` in Visual Studio",
            "Eseguire `ChristmasJumpGame` (profilo Development)",
            "Per generare slide: installare `python-pptx` e lanciare lo script"
        ],
        "notes": "Fornire comandi rapidi per ambiente Windows / PowerShell."
    },
    {
        "title": "Domande e Risorse",
        "bullets": [
            "Repository: percorso locale del progetto",
            "File di riferimento: `BlazorGameEngine-Architecture.md`",
            "Contatti: chiedimi per approfondimenti o demo estese"
        ],
        "notes": "Aprire Q&A. Incoraggiare domande architetturali e su ottimizzazioni."
    }
]


def create_presentation(path_out="BlazorGameEngine_Presentation.pptx"):
    prs = Presentation()

    # Title slide
    first = slides[0]
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = first["title"]
    if len(first["bullets"]) > 0:
        try:
            slide.placeholders[1].text = " - ".join(first["bullets"])[:250]
        except Exception:
            pass
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = first.get("notes", "")

    # Content slides
    for s in slides[1:]:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s["title"]

        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(s["bullets"]):
            p = body.add_paragraph() if i > 0 else body.paragraphs[0]
            p.text = b
            p.level = 0
            p.font.size = Pt(14)

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = s.get("notes", "")

    prs.save(path_out)


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(__file__), "BlazorGameEngine_Presentation.pptx")
    create_presentation(out)
    print(f"Presentazione generata: {out}")
